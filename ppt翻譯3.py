import sys
import os
from PyQt5.QtWidgets import QApplication, QWidget, QVBoxLayout, QPushButton, QFileDialog, QMessageBox
from pptx import Presentation
import openai

openai.api_key = "API KEY"

class PPTTranslatorApp(QWidget):
    def __init__(self):
        super().__init__()
        self.initUI()
    
    def initUI(self):
        self.setWindowTitle('PPT Translator')
        self.setGeometry(100, 100, 300, 100)

        layout = QVBoxLayout()

        self.btn_select_input = QPushButton('Select PPT Files', self)
        self.btn_select_input.clicked.connect(self.select_input_files)
        layout.addWidget(self.btn_select_input)

        self.setLayout(layout)
        self.input_paths = []

    def select_input_files(self):
        options = QFileDialog.Options()
        files, _ = QFileDialog.getOpenFileNames(self, "Select PPT Files", "", "PPT Files (*.pptx);;All Files (*)", options=options)
        if files:
            self.input_paths = files
            desktop_path = os.path.join(os.path.join(os.environ['USERPROFILE']), 'Desktop')
            QMessageBox.information(self, "Files Selected", f"Selected {len(files)} files. Output will be saved to the Desktop.")
            self.translate_pptx()

    def translate_pptx(self):
        if not self.input_paths:
            QMessageBox.warning(self, "Error", "No input files selected.")
            return

        try:
            for file_path in self.input_paths:
                presentation = Presentation(file_path)
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            translated_text = self.translate_text(shape.text)
                            shape.text = translated_text
                
                output_path = os.path.join(
                    os.path.join(os.environ['USERPROFILE'], 'Desktop'),
                    f"translated_{os.path.basename(file_path)}"
                )
                presentation.save(output_path)

            QMessageBox.information(self, "Success", "All translations complete!")
        except Exception as e:
            QMessageBox.critical(self, "Error", str(e))

    def translate_text(self, text, target_language="zh"):
        response = openai.ChatCompletion.create(
            model="gpt-4",
            messages=[
                {"role": "system", "content": "You are a translator."},
                {"role": "user", "content": f"Translate the following text to {target_language}: {text}"}
            ]
        )
        translated_text = response['choices'][0]['message']['content'].strip()
        return translated_text

if __name__ == '__main__':
    app = QApplication(sys.argv)
    translator = PPTTranslatorApp()
    translator.show()
    sys.exit(app.exec_())
